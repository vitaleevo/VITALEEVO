
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Configuration
FILENAME = "VitalEvo_Corporate_Portfolio.pptx"

# Colors
DARK_BG = RGBColor(15, 23, 42)      # #0f172a
PRIMARY = RGBColor(134, 37, 210)    # #8625d2
SECONDARY = RGBColor(16, 185, 129)  # #10b981
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)      # #9ca3af
LIGHT_GRAY = RGBColor(226, 232, 240)# #e2e8f0

# Images
IMG_LOGO = "public/logo.png"
IMG_HERO = "public/hero-card.png"
IMG_BG_LOGIN = "public/login-bg.png"
IMG_OG = "public/og-image.png"

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_header(slide, title, subtitle=None):
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial' # Fallback
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = PRIMARY
        p_sub.font.bold = True

    # Decorative Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY
    line.line.fill.background()

def add_image(slide, img_path, left, top, width=None, height=None):
    try:
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    except Exception as e:
        print(f"Error adding image {img_path}: {e}")

def add_footer(slide, page_num):
    # Footer Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), Inches(13.33), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(50, 50, 70)
    line.line.fill.background()
    
    # Text
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7.35), Inches(5), Inches(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "VitalEvo - Inovação Tecnológica"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    
    # Page Num
    numbox = slide.shapes.add_textbox(Inches(12), Inches(7.35), Inches(1), Inches(0.5))
    tf_num = numbox.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = str(page_num)
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.size = Pt(10)
    p_num.font.color.rgb = GRAY

def add_content_text(slide, text_lines, width=Inches(12)):
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), width, Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for line in text_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(14)

def create_title_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    # Logo at top center (optional) or left
    add_image(slide, IMG_LOGO, Inches(1), Inches(0.5), height=Inches(1))

    # Center content - moved down a bit
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "VITALEVO"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "DESIGN DE ALTO IMPACTO & MARKETING DIGITAL"
    p2.alignment = PP_ALIGN.LEFT
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Conectando Possibilidades em Angola"
    p3.alignment = PP_ALIGN.LEFT
    p3.font.size = Pt(18)
    p3.font.color.rgb = LIGHT_GRAY
    p3.space_before = Pt(10)
    
    # Add Hero Image on the right side
    add_image(slide, IMG_HERO, Inches(7), Inches(1.5), height=Inches(5))
    
    return slide

def create_slide_content(prs, title, subtitle, bullets, page_num, img_path=None):
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    # Small Logo on top right
    add_image(slide, IMG_LOGO, Inches(12), Inches(0.2), height=Inches(0.5))

    add_header(slide, title, subtitle)
    
    # Adjust content width if there is an image
    if img_path:
        add_content_text(slide, bullets, width=Inches(7))
        # Add Image on the right
        add_image(slide, img_path, Inches(8), Inches(2), width=Inches(4.5))
    else:
        add_content_text(slide, bullets)
        
    add_footer(slide, page_num)
    return slide

def main():
    prs = Presentation()
    # Set 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    create_title_slide(prs)

    # 2. Agenda
    create_slide_content(prs, "Agenda", "O que vamos apresentar", [
        "1. Quem Somos",
        "2. Missão, Visão e Valores",
        "3. Nossos Serviços",
        "4. Nossa Metodologia",
        "5. Impacto e Números",
        "6. Compromisso com Angola",
        "7. Futuro e Contato"
    ], 2)

    # 3. Who We Are
    create_slide_content(prs, "Quem Somos", "VitalEvo", [
        "Somos parceiros estratégicos para empresas que desejam liderar o mercado angolano.",
        "Unificamos criatividade, dados e tecnologia para construir marcas fortes e duradouras.",
        "Com sede em Luanda, e operando digitalmente em todas as 18 províncias.",
        "Combinamos design futurista com tecnologia de ponta."
    ], 3, img_path=IMG_OG) # Using OG image for Who We Are

    # 4. Context
    create_slide_content(prs, "O Cenário Atual", "Desafios em Angola", [
        "O mercado angolano está em rápida transformação digital.",
        "Empresas precisam de mais do que apenas um site; precisam de ecossistemas digitais.",
        "A competição exige design premium e performance otimizada.",
        "A segurança de dados e a automação são diferenciais críticos."
    ], 4)

    # 5. Mission
    create_slide_content(prs, "Nossa Missão", "Propósito Central", [
        "Empoderar empresas angolanas com soluções digitais de classe mundial.",
        "Transformar desafios locais em oportunidades globais.",
        "Focar na inovação contínua e na excelência técnica.",
        "Entregar resultados mensuráveis que impulsionam o crescimento real."
    ], 5)

    # 6. Vision
    create_slide_content(prs, "Nossa Visão", "Onde queremos chegar", [
        "Ser o principal catalisador da economia digital em Angola até 2030.",
        "Ser a referência incontestável em qualidade e integridade.",
        "Criar um legado de transformação tecnológica no país."
    ], 6)

    # 7. Values - Part 1
    create_slide_content(prs, "Nossos Valores", "Pilares Fundamentais", [
        "Inovação Constante: Nunca paramos de aprender e evoluir.",
        "Transparência Radical: Honestidade e clareza em todas as etapas.",
        "Foco no Cliente: O sucesso do cliente é o nosso sucesso.",
        "Excelência Angolana: Padrão internacional com DNA local."
    ], 7)

    # 8. Services Overview
    create_slide_content(prs, "Nossos Serviços", "Soluções 360º", [
        "Oferecemos um ecossistema completo de serviços digitais.",
        "• Design Gráfico & UI/UX",
        "• Desenvolvimento Web, Mobile & Software",
        "• Marketing Digital & Gestão de Redes Sociais",
        "• Automação & Inteligência Artificial",
        "• Infraestrutura de TI Corporativa"
    ], 8)
    
    # 9. Service: Design
    create_slide_content(prs, "Design Premium", "Interfaces que Encantam", [
        "Criação de identidades visuais marcantes e memoráveis.",
        "UI/UX Design focado na experiência do usuário e conversão.",
        "Design responsivo para todos os dispositivos e telas.",
        "Estética futurista e alinhada com as tendências globais."
    ], 9)

    # 10. Service: Development
    create_slide_content(prs, "Desenvolvimento High-End", "Web & Mobile", [
        "Desenvolvimento de sites institucionais, e-commerce e portais.",
        "Aplicações Mobile (iOS e Android) nativas e híbridas.",
        "Tecnologias modernas: React, Next.js, Node.js, Python.",
        "Código limpo, escalável e de alta performance."
    ], 10)

    # 11. Service: Marketing
    create_slide_content(prs, "Marketing Digital", "Crescimento Baseado em Dados", [
        "Gestão estratégica de Redes Sociais (Instagram, LinkedIn, Facebook).",
        "Tráfego Pago (Google Ads, Facebook Ads) com foco em ROI.",
        "SEO (Otimização para Motores de Busca) para visibilidade orgânica.",
        "Marketing de Conteúdo e Email Marketing."
    ], 11)

    # 12. Service: AI & Automation
    create_slide_content(prs, "Inteligência Artificial", "O Futuro Agora", [
        "Implementação de Agentes de IA para atendimento e suporte.",
        "Automação de processos repetitivos para aumentar a eficiência.",
        "Análise de dados preditiva para tomada de decisão.",
        "Integração de LLMs (como GPT/Gemini) em fluxos de trabalho."
    ], 12)

    # 13. Service: Infrastructure
    create_slide_content(prs, "Infraestrutura de TI", "Base Sólida", [
        "Consultoria em arquitetura de sistemas e cloud computing.",
        "Segurança da Informação e proteção de dados corporativos.",
        "Manutenção e suporte técnico especializado.",
        "Otimização de bancos de dados e servidores."
    ], 13)

    # 14. Methodology - Overview
    create_slide_content(prs, "Nossa Metodologia", "Processo Comprovado", [
        "Seguimos um processo rigoroso para garantir a excelência.",
        "1. Imersão",
        "2. Estratégia",
        "3. Execução",
        "4. Entrega"
    ], 14)

    # 15. Methodology - Detail 1 & 2
    create_slide_content(prs, "Imersão e Estratégia", "Planejamento", [
        "IMERSÃO: Entendemos profundamente seus objetivos, dores e mercado aprofundando-se na cultura da empresa.",
        "",
        "ESTRATÉGIA: Definimos o plano de ação, tecnologias e cronograma detalhado, alinhando expectativas."
    ], 15)

    # 16. Methodology - Detail 3 & 4
    create_slide_content(prs, "Execução e Entrega", "Ação", [
        "EXECUÇÃO: Nossos designers e desenvolvedores trabalham em sprints ágeis, com feedback constante.",
        "",
        "ENTREGA: Testes rigorosos, lançamento assistido e monitoramento contínuo pós-lançamento."
    ], 16)

    # 17. Stats
    create_slide_content(prs, "Impacto em Números", "Nossa Trajetória", [
        "Mais de 500+ Projetos Entregues com sucesso.",
        "Mais de 250+ Parceiros de Sucesso atendidos.",
        "Mais de 10 Anos de Experiência acumulada no mercado.",
        "Equipe com 45+ Especialistas em diversas áreas."
    ], 17)

    # 18. Commitment to Angola
    create_slide_content(prs, "Compromisso com Angola", "Desenvolvimento Local", [
        "Adaptação de tecnologias globais para a realidade local.",
        "Digitalização de PMEs: A espinha dorsal da economia.",
        "Educação Tecnológica e transferência de conhecimento.",
        "Conectividade Global para produtos angolanos."
    ], 18)

    # 19. Clients & Testimonials
    create_slide_content(prs, "O Que Dizem Nossos Clientes", "Confiança", [
        "\"A VitalEvo transformou nossa presença digital. O novo site triplicou nossos leads.\" - João Silva, CEO Tech Angola",
        "",
        "\"Profissionalismo impecável. A equipe entendeu nossa visão e entregou muito além do esperado.\" - Maria Antónia, Dir. Marketing",
        "",
        "\"Melhor agência em Luanda. Atendimento rápido e alta qualidade.\" - Pedro Kiala, Empreendedor"
    ], 19)

    # 20. Future Roadmap
    create_slide_content(prs, "Visão de Futuro", "2024 - 2030", [
        "Expandir a atuação para toda a região da SADC.",
        "Lançar produtos próprios de SaaS para o mercado africano.",
        "Estabelecer um hub de inovação e aceleração de startups em Luanda.",
        "Continuar liderando a revolução da IA em Angola."
    ], 20)

    # 21. Call to Action
    create_slide_content(prs, "Vamos Começar?", "Transforme seu Negócio", [
        "Estamos prontos para levar sua empresa para o próximo nível.",
        "Junte-se à VitalEvo e faça parte da revolução digital.",
        "",
        "Solicite uma proposta personalizada hoje mesmo."
    ], 21, img_path=IMG_BG_LOGIN) # Using Login BG as highlight

    # 22. Contact
    create_slide_content(prs, "Contatos", "Fale Conosco", [
        "Website: www.vitaleevo.ao",
        "Email: contato@vitaleevo.ao",
        "Telefone: +244 923 456 789",
        "Endereço: Luanda, Angola",
        "",
        "Siga-nos nas redes sociais: @vitaleevo"
    ], 22, img_path=IMG_LOGO)

    prs.save(FILENAME)
    print(f"Presentation saved to {FILENAME}")

if __name__ == "__main__":
    main()
